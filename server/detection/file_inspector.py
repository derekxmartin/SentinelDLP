"""File content inspector — extracts text from common file formats.

Decomposes files into ParsedMessage components for downstream detection.
Supports: PDF, DOCX, XLSX, PPTX, plain text (with encoding detection),
EML (with recursive attachment extraction), CSV, HTML, and RTF.

Each extracted component preserves metadata about its origin (page number,
sheet name, slide number, attachment filename, etc.).
"""

from __future__ import annotations

import csv
import email
import email.policy
import io
import logging
from email.message import EmailMessage
from html.parser import HTMLParser
from pathlib import Path

import chardet

from server.detection.models import ComponentType, ParsedMessage

logger = logging.getLogger(__name__)


# Maximum file size for extraction (100MB)
MAX_FILE_SIZE = 100 * 1024 * 1024

# Maximum text length per component to prevent memory issues
MAX_COMPONENT_TEXT = 10 * 1024 * 1024  # 10MB


class _HTMLTextExtractor(HTMLParser):
    """Simple HTML-to-text extractor that strips tags."""

    def __init__(self):
        super().__init__()
        self._text: list[str] = []
        self._skip = False

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip = True

    def handle_endtag(self, tag):
        if tag in ("script", "style"):
            self._skip = False
        if tag in ("p", "br", "div", "h1", "h2", "h3", "h4", "h5", "h6", "li", "tr"):
            self._text.append("\n")

    def handle_data(self, data):
        if not self._skip:
            self._text.append(data)

    def get_text(self) -> str:
        return "".join(self._text).strip()


def _truncate(text: str) -> str:
    """Truncate text to MAX_COMPONENT_TEXT."""
    if len(text) > MAX_COMPONENT_TEXT:
        return text[:MAX_COMPONENT_TEXT] + "\n[TRUNCATED]"
    return text


class FileInspector:
    """Extracts text content from files and returns ParsedMessage.

    Usage:
        >>> inspector = FileInspector()
        >>> message = inspector.inspect(file_bytes, filename="report.pdf")
        >>> for comp in message.components:
        ...     print(comp.component_type, comp.name, len(comp.content))

    Supported formats:
        - PDF (.pdf) via pdfplumber
        - DOCX (.docx) via python-docx
        - XLSX (.xlsx) via openpyxl
        - PPTX (.pptx) via python-pptx
        - Plain text (.txt, .log, .md, .csv, etc.) with chardet encoding detection
        - EML (.eml) via email stdlib — decomposes envelope, subject, body, attachments
        - HTML (.html, .htm) — strips tags
    """

    # Extension to handler method mapping
    _HANDLERS: dict[str, str] = {
        ".pdf": "_extract_pdf",
        ".docx": "_extract_docx",
        ".xlsx": "_extract_xlsx",
        ".pptx": "_extract_pptx",
        ".eml": "_extract_eml",
        ".html": "_extract_html",
        ".htm": "_extract_html",
    }

    # Extensions treated as plain text
    _TEXT_EXTENSIONS: set[str] = {
        ".txt", ".log", ".md", ".csv", ".tsv", ".json", ".xml",
        ".yaml", ".yml", ".ini", ".cfg", ".conf", ".py", ".js",
        ".ts", ".java", ".c", ".cpp", ".h", ".hpp", ".cs", ".go",
        ".rs", ".rb", ".php", ".sh", ".bat", ".ps1", ".sql",
        ".rtf",
    }

    def inspect(
        self,
        content: bytes,
        filename: str = "",
        metadata: dict | None = None,
    ) -> ParsedMessage:
        """Extract text content from a file and return a ParsedMessage.

        Args:
            content: Raw file bytes.
            filename: Original filename (used for format detection).
            metadata: Optional metadata to attach to the message
                (e.g., sender, channel).

        Returns:
            ParsedMessage with extracted text components.
        """
        if len(content) > MAX_FILE_SIZE:
            logger.warning(
                "File %s exceeds max size (%d > %d), skipping",
                filename,
                len(content),
                MAX_FILE_SIZE,
            )
            msg = ParsedMessage(metadata=metadata or {})
            msg.add_component(
                ComponentType.GENERIC,
                f"[File too large: {len(content)} bytes]",
                {"filename": filename, "error": "exceeds_max_size"},
            )
            return msg

        ext = ""
        if filename:
            ext = Path(filename).suffix.lower()

        msg = ParsedMessage(metadata=metadata or {})

        # Route to appropriate handler
        handler_name = self._HANDLERS.get(ext)
        if handler_name:
            handler = getattr(self, handler_name)
            try:
                handler(content, filename, msg)
            except Exception as exc:
                logger.error("Failed to extract %s: %s", filename, exc)
                msg.add_component(
                    ComponentType.GENERIC,
                    f"[Extraction failed: {exc}]",
                    {"filename": filename, "error": str(exc)},
                )
        elif ext in self._TEXT_EXTENSIONS or self._looks_like_text(content):
            self._extract_text(content, filename, msg)
        else:
            # Unknown binary — store filename metadata only
            msg.add_component(
                ComponentType.ATTACHMENT,
                "",
                {
                    "filename": filename,
                    "size": len(content),
                    "content_bytes": content,
                    "note": "unsupported_format",
                },
            )

        return msg

    def _extract_pdf(
        self, content: bytes, filename: str, msg: ParsedMessage
    ) -> None:
        """Extract text from PDF using pdfplumber."""
        import pdfplumber

        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                # Also extract table text
                tables = page.extract_tables() or []
                table_text = ""
                for table in tables:
                    for row in table:
                        if row:
                            cells = [str(c) if c is not None else "" for c in row]
                            table_text += " | ".join(cells) + "\n"

                combined = text
                if table_text:
                    combined += "\n[TABLE]\n" + table_text

                if combined.strip():
                    msg.add_component(
                        ComponentType.BODY,
                        _truncate(combined),
                        {
                            "filename": filename,
                            "page": i,
                            "total_pages": len(pdf.pages),
                            "source": "pdf",
                        },
                    )

    def _extract_docx(
        self, content: bytes, filename: str, msg: ParsedMessage
    ) -> None:
        """Extract text from DOCX using python-docx."""
        from docx import Document

        doc = Document(io.BytesIO(content))

        # Extract paragraphs
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        if paragraphs:
            msg.add_component(
                ComponentType.BODY,
                _truncate("\n".join(paragraphs)),
                {"filename": filename, "source": "docx_paragraphs"},
            )

        # Extract tables
        for i, table in enumerate(doc.tables, 1):
            rows = []
            for row in table.rows:
                cells = [cell.text for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                msg.add_component(
                    ComponentType.BODY,
                    _truncate("\n".join(rows)),
                    {
                        "filename": filename,
                        "table_index": i,
                        "source": "docx_table",
                    },
                )

    def _extract_xlsx(
        self, content: bytes, filename: str, msg: ParsedMessage
    ) -> None:
        """Extract text from XLSX using openpyxl."""
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))

            if rows:
                msg.add_component(
                    ComponentType.BODY,
                    _truncate("\n".join(rows)),
                    {
                        "filename": filename,
                        "sheet": sheet_name,
                        "source": "xlsx",
                    },
                )

        wb.close()

    def _extract_pptx(
        self, content: bytes, filename: str, msg: ParsedMessage
    ) -> None:
        """Extract text from PPTX using python-pptx."""
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))

        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text for cell in row.cells]
                        texts.append(" | ".join(cells))

            if texts:
                msg.add_component(
                    ComponentType.BODY,
                    _truncate("\n".join(texts)),
                    {
                        "filename": filename,
                        "slide": i,
                        "total_slides": len(prs.slides),
                        "source": "pptx",
                    },
                )

    def _extract_eml(
        self, content: bytes, filename: str, msg: ParsedMessage
    ) -> None:
        """Extract components from EML using email stdlib.

        Decomposes into: envelope (from/to), subject, body, attachments.
        """
        eml = email.message_from_bytes(content, policy=email.policy.default)

        # Envelope
        envelope_parts = []
        for header in ("From", "To", "Cc", "Bcc", "Date", "Message-ID"):
            val = eml.get(header)
            if val:
                envelope_parts.append(f"{header}: {val}")
        if envelope_parts:
            msg.add_component(
                ComponentType.ENVELOPE,
                "\n".join(envelope_parts),
                {"filename": filename, "source": "eml_envelope"},
            )

        # Subject
        subject = eml.get("Subject", "")
        if subject:
            msg.add_component(
                ComponentType.SUBJECT,
                subject,
                {"filename": filename, "source": "eml_subject"},
            )

        # Body
        body = eml.get_body(preferencelist=("plain", "html"))
        if body:
            body_content = body.get_content()
            if isinstance(body_content, bytes):
                body_content = body_content.decode("utf-8", errors="replace")

            # If HTML body, strip tags
            content_type = body.get_content_type()
            if content_type == "text/html":
                extractor = _HTMLTextExtractor()
                extractor.feed(body_content)
                body_content = extractor.get_text()

            if body_content.strip():
                msg.add_component(
                    ComponentType.BODY,
                    _truncate(body_content),
                    {
                        "filename": filename,
                        "content_type": content_type,
                        "source": "eml_body",
                    },
                )

        # Attachments
        for part in eml.iter_attachments():
            att_filename = part.get_filename() or "unnamed_attachment"
            att_content = part.get_content()

            if isinstance(att_content, str):
                att_bytes = att_content.encode("utf-8")
                att_text = att_content
            elif isinstance(att_content, bytes):
                att_bytes = att_content
                # Try to extract text from the attachment recursively
                sub_msg = self.inspect(att_bytes, att_filename)
                att_text = "\n".join(
                    c.content for c in sub_msg.components if c.content
                )
            else:
                att_bytes = b""
                att_text = ""

            msg.add_component(
                ComponentType.ATTACHMENT,
                _truncate(att_text) if att_text else "",
                {
                    "filename": att_filename,
                    "size": len(att_bytes),
                    "content_bytes": att_bytes,
                    "content_type": part.get_content_type(),
                    "source": "eml_attachment",
                },
            )

    def _extract_html(
        self, content: bytes, filename: str, msg: ParsedMessage
    ) -> None:
        """Extract text from HTML by stripping tags."""
        encoding = self._detect_encoding(content)
        text = content.decode(encoding, errors="replace")

        extractor = _HTMLTextExtractor()
        extractor.feed(text)
        plain = extractor.get_text()

        if plain.strip():
            msg.add_component(
                ComponentType.BODY,
                _truncate(plain),
                {
                    "filename": filename,
                    "encoding": encoding,
                    "source": "html",
                },
            )

    def _extract_text(
        self, content: bytes, filename: str, msg: ParsedMessage
    ) -> None:
        """Extract plain text with encoding detection."""
        encoding = self._detect_encoding(content)
        text = content.decode(encoding, errors="replace")

        if text.strip():
            msg.add_component(
                ComponentType.BODY,
                _truncate(text),
                {
                    "filename": filename,
                    "encoding": encoding,
                    "source": "text",
                },
            )

    @staticmethod
    def _detect_encoding(content: bytes) -> str:
        """Detect text encoding using chardet.

        Falls back to utf-8 if detection confidence is low.
        """
        if not content:
            return "utf-8"

        result = chardet.detect(content)
        encoding = result.get("encoding") or "utf-8"
        confidence = result.get("confidence", 0)

        # Fall back to utf-8 if confidence is too low
        if confidence < 0.5:
            encoding = "utf-8"

        # Normalize encoding names
        encoding = encoding.lower().replace("-", "_")

        logger.debug(
            "Detected encoding: %s (confidence: %.2f)",
            encoding,
            confidence,
        )
        return encoding

    @staticmethod
    def _looks_like_text(content: bytes) -> bool:
        """Heuristic: check if content appears to be text.

        Samples the first 8KB and checks if it's mostly printable.
        """
        sample = content[:8192]
        if not sample:
            return False

        # Count non-text bytes (excluding common whitespace)
        text_bytes = set(range(0x20, 0x7F)) | {0x09, 0x0A, 0x0D}
        non_text = sum(1 for b in sample if b not in text_bytes and b < 0x80)

        # Allow up to 5% non-text bytes (for UTF-8 multibyte)
        return non_text / len(sample) < 0.05
