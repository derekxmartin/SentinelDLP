"""Tests for FileInspector (P1-T6).

Covers: PDF extraction, DOCX with tables, XLSX multi-sheet, PPTX,
plain text with encoding detection (Shift-JIS), EML with attachment,
HTML stripping, error handling, and edge cases.
"""

import io
from email.message import EmailMessage


from server.detection.models import ComponentType
from server.detection.file_inspector import FileInspector


# ---------------------------------------------------------------------------
# Helpers — synthetic file generators
# ---------------------------------------------------------------------------


def _make_pdf_with_text(*pages: str) -> bytes:
    """Create a real PDF with text content using pdfplumber-compatible format.

    Uses a minimal but valid PDF structure.
    """
    # Use reportlab-free approach: build minimal PDF by hand
    obj_num = 1

    # Catalog
    catalog_num = obj_num
    obj_num += 1

    # Pages
    pages_num = obj_num
    obj_num += 1

    # Font
    font_num = obj_num
    obj_num += 1

    page_objects = []
    for page_text in pages:
        page_num = obj_num
        obj_num += 1
        content_num = obj_num
        obj_num += 1
        page_objects.append((page_num, content_num, page_text))

    # Build PDF
    pdf = io.BytesIO()
    pdf.write(b"%PDF-1.4\n")
    offsets = {}

    # Font object
    offsets[font_num] = pdf.tell()
    pdf.write(f"{font_num} 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n".encode())

    # Content streams and page objects
    kids = []
    for page_num, content_num, text in page_objects:
        # Content stream
        stream_content = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode()
        offsets[content_num] = pdf.tell()
        pdf.write(f"{content_num} 0 obj\n<< /Length {len(stream_content)} >>\nstream\n".encode())
        pdf.write(stream_content)
        pdf.write(b"\nendstream\nendobj\n")

        # Page object
        offsets[page_num] = pdf.tell()
        pdf.write(
            f"{page_num} 0 obj\n"
            f"<< /Type /Page /Parent {pages_num} 0 R "
            f"/MediaBox [0 0 612 792] "
            f"/Contents {content_num} 0 R "
            f"/Resources << /Font << /F1 {font_num} 0 R >> >> "
            f">>\nendobj\n".encode()
        )
        kids.append(f"{page_num} 0 R")

    # Pages object
    offsets[pages_num] = pdf.tell()
    kids_str = " ".join(kids)
    pdf.write(
        f"{pages_num} 0 obj\n"
        f"<< /Type /Pages /Kids [{kids_str}] /Count {len(pages)} >>\n"
        f"endobj\n".encode()
    )

    # Catalog
    offsets[catalog_num] = pdf.tell()
    pdf.write(
        f"{catalog_num} 0 obj\n"
        f"<< /Type /Catalog /Pages {pages_num} 0 R >>\n"
        f"endobj\n".encode()
    )

    # Cross-reference table
    xref_offset = pdf.tell()
    pdf.write(b"xref\n")
    pdf.write(f"0 {obj_num}\n".encode())
    pdf.write(b"0000000000 65535 f \n")
    for i in range(1, obj_num):
        pdf.write(f"{offsets.get(i, 0):010d} 00000 n \n".encode())

    # Trailer
    pdf.write(
        f"trailer\n<< /Size {obj_num} /Root {catalog_num} 0 R >>\n"
        f"startxref\n{xref_offset}\n%%EOF\n".encode()
    )

    return pdf.getvalue()


def _make_docx_with_content(paragraphs: list[str], tables: list[list[list[str]]] | None = None) -> bytes:
    """Create a DOCX with paragraphs and optional tables."""
    from docx import Document

    doc = Document()
    for para in paragraphs:
        doc.add_paragraph(para)

    if tables:
        for table_data in tables:
            if not table_data:
                continue
            cols = len(table_data[0])
            table = doc.add_table(rows=len(table_data), cols=cols)
            for i, row_data in enumerate(table_data):
                for j, cell_text in enumerate(row_data):
                    table.rows[i].cells[j].text = cell_text

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx_with_sheets(sheets: dict[str, list[list]]) -> bytes:
    """Create an XLSX with named sheets and data."""
    from openpyxl import Workbook

    wb = Workbook()
    # Remove default sheet
    wb.remove(wb.active)

    for name, rows in sheets.items():
        ws = wb.create_sheet(title=name)
        for row in rows:
            ws.append(row)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx_with_slides(slides: list[str]) -> bytes:
    """Create a PPTX with text on each slide."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for text in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        txBox.text_frame.text = text

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_eml(
    from_addr: str = "sender@test.com",
    to_addr: str = "recipient@test.com",
    subject: str = "Test Email",
    body: str = "This is the email body.",
    attachments: list[tuple[str, bytes, str]] | None = None,
) -> bytes:
    """Create an EML file.

    attachments: list of (filename, content_bytes, mime_type)
    """
    msg = EmailMessage()
    msg["From"] = from_addr
    msg["To"] = to_addr
    msg["Subject"] = subject
    msg.set_content(body)

    if attachments:
        for att_name, att_content, att_mime in attachments:
            maintype, subtype = att_mime.split("/", 1)
            msg.add_attachment(
                att_content,
                maintype=maintype,
                subtype=subtype,
                filename=att_name,
            )

    return msg.as_bytes()


# ---------------------------------------------------------------------------
# PDF extraction tests
# ---------------------------------------------------------------------------


class TestPDFExtraction:

    def test_single_page_pdf(self):
        inspector = FileInspector()
        pdf = _make_pdf_with_text("Hello World from page one")
        msg = inspector.inspect(pdf, "test.pdf")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert len(body_comps) >= 1
        assert "Hello World" in body_comps[0].content

    def test_3_page_pdf(self):
        """Acceptance: extract text from 3-page PDF."""
        inspector = FileInspector()
        pdf = _make_pdf_with_text(
            "Page one: confidential data",
            "Page two: credit card 4532015112830366",
            "Page three: SSN 123-45-6789",
        )
        msg = inspector.inspect(pdf, "report.pdf")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert len(body_comps) == 3

        all_text = " ".join(c.content for c in body_comps)
        assert "confidential" in all_text
        assert "4532015112830366" in all_text
        assert "123-45-6789" in all_text

    def test_pdf_page_metadata(self):
        inspector = FileInspector()
        pdf = _make_pdf_with_text("Page 1", "Page 2")
        msg = inspector.inspect(pdf, "doc.pdf")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert body_comps[0].metadata["page"] == 1
        assert body_comps[1].metadata["page"] == 2
        assert body_comps[0].metadata["total_pages"] == 2


# ---------------------------------------------------------------------------
# DOCX extraction tests
# ---------------------------------------------------------------------------


class TestDOCXExtraction:

    def test_paragraphs(self):
        inspector = FileInspector()
        docx = _make_docx_with_content([
            "First paragraph with sensitive data.",
            "Second paragraph with SSN 123-45-6789.",
        ])
        msg = inspector.inspect(docx, "doc.docx")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert len(body_comps) >= 1
        all_text = " ".join(c.content for c in body_comps)
        assert "sensitive data" in all_text
        assert "123-45-6789" in all_text

    def test_docx_with_tables(self):
        """Acceptance: DOCX with tables."""
        inspector = FileInspector()
        docx = _make_docx_with_content(
            ["Document with table below:"],
            tables=[[
                ["Name", "SSN", "Salary"],
                ["Alice", "123-45-6789", "$100,000"],
                ["Bob", "987-65-4321", "$95,000"],
            ]],
        )
        msg = inspector.inspect(docx, "employees.docx")

        all_text = " ".join(c.content for c in msg.components)
        assert "Alice" in all_text
        assert "123-45-6789" in all_text
        assert "Bob" in all_text

    def test_docx_metadata(self):
        inspector = FileInspector()
        docx = _make_docx_with_content(["Test"])
        msg = inspector.inspect(docx, "test.docx")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert body_comps[0].metadata["source"] == "docx_paragraphs"


# ---------------------------------------------------------------------------
# XLSX extraction tests
# ---------------------------------------------------------------------------


class TestXLSXExtraction:

    def test_single_sheet(self):
        inspector = FileInspector()
        xlsx = _make_xlsx_with_sheets({
            "Data": [
                ["Name", "Email"],
                ["Alice", "alice@example.com"],
            ]
        })
        msg = inspector.inspect(xlsx, "data.xlsx")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert len(body_comps) == 1
        assert "alice@example.com" in body_comps[0].content

    def test_multiple_sheets(self):
        """Acceptance: XLSX with multiple sheets."""
        inspector = FileInspector()
        xlsx = _make_xlsx_with_sheets({
            "Employees": [
                ["Name", "SSN"],
                ["Alice", "123-45-6789"],
            ],
            "Financial": [
                ["Account", "Balance"],
                ["Checking", "$50,000"],
            ],
            "Cards": [
                ["Type", "Number"],
                ["Visa", "4532015112830366"],
            ],
        })
        msg = inspector.inspect(xlsx, "company.xlsx")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert len(body_comps) == 3

        sheets = {c.metadata["sheet"] for c in body_comps}
        assert sheets == {"Employees", "Financial", "Cards"}

        all_text = " ".join(c.content for c in body_comps)
        assert "123-45-6789" in all_text
        assert "4532015112830366" in all_text

    def test_xlsx_sheet_metadata(self):
        inspector = FileInspector()
        xlsx = _make_xlsx_with_sheets({"Sheet1": [["data"]]})
        msg = inspector.inspect(xlsx, "test.xlsx")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert body_comps[0].metadata["sheet"] == "Sheet1"
        assert body_comps[0].metadata["source"] == "xlsx"


# ---------------------------------------------------------------------------
# PPTX extraction tests
# ---------------------------------------------------------------------------


class TestPPTXExtraction:

    def test_slides(self):
        inspector = FileInspector()
        pptx = _make_pptx_with_slides([
            "Slide 1: Quarterly Report",
            "Slide 2: Revenue $1.2M",
            "Slide 3: Confidential projections",
        ])
        msg = inspector.inspect(pptx, "presentation.pptx")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert len(body_comps) == 3

        all_text = " ".join(c.content for c in body_comps)
        assert "Quarterly Report" in all_text
        assert "Revenue" in all_text
        assert "Confidential" in all_text

    def test_pptx_metadata(self):
        inspector = FileInspector()
        pptx = _make_pptx_with_slides(["Slide 1"])
        msg = inspector.inspect(pptx, "deck.pptx")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert body_comps[0].metadata["slide"] == 1
        assert body_comps[0].metadata["source"] == "pptx"


# ---------------------------------------------------------------------------
# Plain text + encoding detection
# ---------------------------------------------------------------------------


class TestTextExtraction:

    def test_utf8(self):
        inspector = FileInspector()
        content = "Hello UTF-8 world with SSN 123-45-6789".encode("utf-8")
        msg = inspector.inspect(content, "data.txt")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert len(body_comps) == 1
        assert "123-45-6789" in body_comps[0].content

    def test_shift_jis_detected(self):
        """Acceptance: Shift-JIS text detected and decoded."""
        inspector = FileInspector()
        # Longer Japanese text for reliable chardet detection
        japanese_text = (
            "これは機密文書です。社会保障番号：123-45-6789\n"
            "この文書は厳重に管理されています。\n"
            "従業員の個人情報が含まれています。\n"
            "クレジットカード番号や銀行口座情報。\n"
            "取り扱いには十分注意してください。\n"
            "不正アクセスは法律で禁止されています。\n"
        )
        content = japanese_text.encode("shift_jis")
        msg = inspector.inspect(content, "document.txt")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert len(body_comps) == 1
        # Content should contain the SSN regardless of detected encoding name
        assert "123-45-6789" in body_comps[0].content
        # Encoding metadata should be present
        assert "encoding" in body_comps[0].metadata

    def test_latin1(self):
        inspector = FileInspector()
        content = "Données personnelles: crédit".encode("latin-1")
        msg = inspector.inspect(content, "french.txt")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert len(body_comps) == 1

    def test_csv_as_text(self):
        inspector = FileInspector()
        content = b"Name,SSN,Email\nAlice,123-45-6789,alice@test.com\n"
        msg = inspector.inspect(content, "data.csv")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert len(body_comps) == 1
        assert "123-45-6789" in body_comps[0].content

    def test_source_code(self):
        inspector = FileInspector()
        content = b'API_KEY = "sk-secret-key-12345"\nDATABASE_URL = "postgres://user:pass@localhost/db"'
        msg = inspector.inspect(content, "config.py")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert len(body_comps) == 1
        assert "sk-secret-key" in body_comps[0].content


# ---------------------------------------------------------------------------
# EML extraction
# ---------------------------------------------------------------------------


class TestEMLExtraction:

    def test_basic_email(self):
        inspector = FileInspector()
        eml = _make_eml(
            subject="Sensitive Report",
            body="The SSN is 123-45-6789 and the CC is 4532015112830366.",
        )
        msg = inspector.inspect(eml, "email.eml")

        # Should have envelope, subject, body
        comp_types = {c.component_type for c in msg.components}
        assert ComponentType.ENVELOPE in comp_types
        assert ComponentType.SUBJECT in comp_types
        assert ComponentType.BODY in comp_types

        # Check content
        subject_comp = next(c for c in msg.components if c.component_type == ComponentType.SUBJECT)
        assert subject_comp.content == "Sensitive Report"

        body_comp = next(c for c in msg.components if c.component_type == ComponentType.BODY)
        assert "123-45-6789" in body_comp.content

    def test_eml_with_attachment(self):
        """Acceptance: EML with attachment extracted."""
        inspector = FileInspector()
        att_content = b"Name,SSN\nAlice,123-45-6789\n"
        eml = _make_eml(
            subject="Data Export",
            body="Please find the data attached.",
            attachments=[("employees.csv", att_content, "text/csv")],
        )
        msg = inspector.inspect(eml, "export.eml")

        # Should have attachment component
        att_comps = [c for c in msg.components if c.component_type == ComponentType.ATTACHMENT]
        assert len(att_comps) == 1
        assert att_comps[0].metadata["filename"] == "employees.csv"

    def test_eml_envelope_headers(self):
        inspector = FileInspector()
        eml = _make_eml(
            from_addr="ceo@company.com",
            to_addr="all@company.com",
            subject="Quarterly Results",
            body="Revenue figures attached.",
        )
        msg = inspector.inspect(eml, "email.eml")

        envelope = next(
            (c for c in msg.components if c.component_type == ComponentType.ENVELOPE),
            None,
        )
        assert envelope is not None
        assert "ceo@company.com" in envelope.content
        assert "all@company.com" in envelope.content


# ---------------------------------------------------------------------------
# HTML extraction
# ---------------------------------------------------------------------------


class TestHTMLExtraction:

    def test_html_strips_tags(self):
        inspector = FileInspector()
        html = b"<html><body><h1>Title</h1><p>SSN: 123-45-6789</p></body></html>"
        msg = inspector.inspect(html, "page.html")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert len(body_comps) == 1
        assert "123-45-6789" in body_comps[0].content
        assert "<p>" not in body_comps[0].content

    def test_html_strips_script_and_style(self):
        inspector = FileInspector()
        html = b"<html><head><style>body{color:red}</style></head><body><script>alert(1)</script><p>Secret data</p></body></html>"
        msg = inspector.inspect(html, "page.html")

        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        text = body_comps[0].content
        assert "Secret data" in text
        assert "alert" not in text
        assert "color:red" not in text


# ---------------------------------------------------------------------------
# Edge cases
# ---------------------------------------------------------------------------


class TestEdgeCases:

    def test_empty_file(self):
        inspector = FileInspector()
        msg = inspector.inspect(b"", "empty.txt")
        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert len(body_comps) == 0

    def test_unknown_binary(self):
        inspector = FileInspector()
        msg = inspector.inspect(b"\x00\x01\x02\x03" * 100, "unknown.bin")
        # Should store as attachment with metadata
        att_comps = [c for c in msg.components if c.component_type == ComponentType.ATTACHMENT]
        assert len(att_comps) == 1
        assert att_comps[0].metadata.get("note") == "unsupported_format"

    def test_corrupt_pdf_graceful(self):
        inspector = FileInspector()
        msg = inspector.inspect(b"%PDF-1.4 corrupt content", "bad.pdf")
        # Should have an error component, not crash
        assert len(msg.components) >= 1

    def test_metadata_preserved(self):
        inspector = FileInspector()
        msg = inspector.inspect(
            b"test content",
            "test.txt",
            metadata={"sender": "user@test.com", "channel": "email"},
        )
        assert msg.metadata["sender"] == "user@test.com"
        assert msg.metadata["channel"] == "email"

    def test_no_filename(self):
        """File with no name — uses text heuristic."""
        inspector = FileInspector()
        msg = inspector.inspect(b"Plain text without a filename")
        body_comps = [c for c in msg.components if c.component_type == ComponentType.BODY]
        assert len(body_comps) == 1


# ---------------------------------------------------------------------------
# Integration: FileInspector → DetectionEngine
# ---------------------------------------------------------------------------


class TestIntegrationWithEngine:

    def test_inspect_then_detect(self):
        """Full pipeline: inspect file → run detection engine."""
        from server.detection.engine import DetectionEngine
        from server.detection.analyzers.regex_analyzer import RegexAnalyzer, RegexPattern

        inspector = FileInspector()
        engine = DetectionEngine()
        engine.register(
            RegexAnalyzer(
                name="ssn",
                patterns=[
                    RegexPattern(
                        name="US SSN",
                        pattern=r"\b\d{3}-\d{2}-\d{4}\b",
                    )
                ],
            )
        )

        # Inspect a DOCX → get ParsedMessage → run detection
        docx = _make_docx_with_content([
            "Employee record: SSN 123-45-6789",
            "Another record: SSN 987-65-4321",
        ])
        msg = inspector.inspect(docx, "employees.docx")
        result = engine.detect(msg)

        assert result.has_matches
        assert result.match_count >= 2
        ssn_matches = [m for m in result.matches if m.rule_name == "US SSN"]
        texts = {m.matched_text for m in ssn_matches}
        assert "123-45-6789" in texts
        assert "987-65-4321" in texts
